from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.exc import PackageNotFoundError
import sys


def get_title_object(input_presentation):
    main_obj = dict()
    try :
        prs = Presentation(input_presentation)
        for index, slide in enumerate(prs.slides):
            for indexSlide, shape in enumerate(slide.shapes):
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == PP_PLACEHOLDER.TITLE or phf.type == PP_PLACEHOLDER.CENTER_TITLE or phf.type == PP_PLACEHOLDER.VERTICAL_TITLE:
                        main_obj[index] = shape.text_frame.text
        return main_obj
    except PackageNotFoundError:
        print("PPTX not Found.")


def main():
    try :
        str(sys.argv[1])
        print get_title_object(str(sys.argv[1]))
    except IndexError:
        print("File Argument not given.")


if __name__ == "__main__": main()
